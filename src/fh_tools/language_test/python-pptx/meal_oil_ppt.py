# -*- coding: utf-8 -*-
"""
Created on Mon Mar  5 15:07:53 2018

@author: 杨宝元
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Pt,Inches
prs = Presentation()
import os
from utils.analist_tools import fetch_edb
from utils.analist_tools import period_plot,plot,twin_plot
from backend import  engine_md
import pandas as pd
import numpy as np
from scipy.interpolate import interp1d
import datetime
import matplotlib.pyplot as plt
from pylab import mpl
mpl.rcParams['font.sans-serif'] = ['SimHei']  # 字体可以根据需要改动
mpl.rcParams['axes.unicode_minus'] = False  # 解决中文减号不显示的问题

import datetime


# prs = Presentation()
# slide = prs.slides.add_slide(prs.slide_layouts[0])
# body_shape = slide.shapes.placeholders  # body_shape为本页ppt中所有shapes
# body_shape[0].text = '油脂油料行情分析报告'  # 在第一个文本框中文字框架内添加文字
# body_shape[1].text = '中国国际期货期货 \n 研究院 \n %s '% datetime.datetime.now().strftime("%Y-%m-%d")  # 在第二个文本框中文字框架内添加文字
# img_path = 'logo.png'  # 文件路径
# # left, top, width, height = Inches(7), Inches(0), Inches(8), Inches(5)  # 预设位置及大小
# pic = slide.shapes.add_picture(img_path, left=Inches(8),top=Inches(0.2), width=Inches(1.5), height=Inches(0.35))  # 在指定位置按预设值添加图片
# # 对ppt的修改
# # path_oil_meal=r'D:\future\oil_meal'
# #保存ppt文档
def add_ppt(title,text_input,img_path,value1,value2,value3):
    # prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    log_path = 'logo.png'  # 文件路径
    pic1 = slide.shapes.add_picture(log_path, left=Inches(8), top=Inches(0.2), width=Inches(1.5),
                                   height=Inches(0.35))  # 在指定位置按预设值添加图片
    txBox = slide.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    # 加新的段落
    p = tf.add_paragraph()
    p.text = text_input.format(value1,value2,value3)
    p.font.size = Pt(24)
    left, top, width, height = Inches(1), Inches(2.3), Inches(8), Inches(5)  # 预设位置及大小
    pic2 = slide.shapes.add_picture(img_path, left, top, width, height)  # 在指定位置按预设值添加图片
#养殖利润
def add_ppt_oil_meal():
    path_oil_meal = os.path.join(os.path.abspath(os.path.curdir), 'output', 'oil_meal')
    if not os.path.exists(path_oil_meal):
        os.makedirs(path_oil_meal)
    sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM oil_meal_edb edb left join oil_meal_info info 
    on edb.key=info.key where en_name like 'farming_profit%%'"""
    farming_profit=fetch_edb(sql_str)
    farming_profit.tail()
    path=path_oil_meal+'\\自繁养猪利润.jpg'
    period_plot(list(['2013','2014','2015','2016','2017']),farming_profit.farming_profit_homebred_pig.resample('W-FRI').mean().interpolate(),'2018','元/头','自繁养猪利润',path)
    add_ppt(title='自繁养猪利润',text_input='本期自繁生猪养殖利润为{}元/头,上期养殖利润{}元/头,\n变动{}元/头,整体利润表现',img_path=path_oil_meal+'\\自繁养猪利润.jpg',
            value1=farming_profit.farming_profit_homebred_pig[-1],value2=farming_profit.farming_profit_homebred_pig[-2],
            value3=round(farming_profit.farming_profit_homebred_pig[-1]-farming_profit.farming_profit_homebred_pig[-2],2))
    path=path_oil_meal+'\\肉鸭养殖利润.jpg'
    period_plot(list(['2013','2014','2015','2016','2017']),farming_profit.farming_profit_duck.resample('W-FRI').mean().interpolate(),'2018','元/羽','肉鸭养殖利润',path)
    add_ppt(title='肉鸭养殖利润',text_input='本期肉鸭养殖利润为{}元/羽,上期养殖利润{}元/羽,\n变动{}元/羽,整体利润表现',img_path=path_oil_meal+'\\肉鸭养殖利润.jpg',
            value1=farming_profit.farming_profit_duck[-1],value2=farming_profit.farming_profit_duck[-2],
            value3=round(farming_profit.farming_profit_duck[-1]-farming_profit.farming_profit_duck[-2],2))
    path=path_oil_meal+'\\蛋鸡养殖利润.jpg'
    period_plot(list(['2013','2014','2015','2016','2017']),farming_profit.farming_profit_layer.resample('W-FRI').mean().interpolate(),'2018','元/羽','蛋鸡养殖利润',path)
    add_ppt(title='蛋鸡养殖利润',text_input='本期蛋鸡养殖利润为{}元/羽,上期养殖利润{}元/羽,\n变动{}元/羽,整体利润表现',img_path=path_oil_meal+'\\蛋鸡养殖利润.jpg',
            value1=farming_profit.farming_profit_layer[-1],value2=farming_profit.farming_profit_layer[-2],
            value3=round(farming_profit.farming_profit_layer[-1]-farming_profit.farming_profit_layer[-2],2))

    prs.save(path_oil_meal+'\\python-pptx.pptx')

# #油脂国际价格
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('peanut_oil_price_international','soybean_oil_price_international','cottonseed_oil_price_international',
# 'corn_oil_price_international','palm_oil_price_international')"""
# oil_price_international= fetch_edb(sql_str,language='cn')
# path=path_oil_meal+'\\国际油脂现货价格.jpg'
# plot(oil_price_international['2016':].resample('W').mean().interpolate(),'时间','美分/磅','国际油脂现货价格',path,marker='')
# #豆粕现货价格
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('soybean_meal_price_Dongguan','soybean_meal_price_Lianyungang','soybean_meal_price_Rizhao'
# 'soybean_meal_price_Fangchenggang','soybean_meal_price_Tianjing','soybean_meal_price_Chengdu','soybean_meal_price_moa','soybean_meal_price_Zhangjiagang','soybean_meal_price_Rizhao') """
# m_price= fetch_edb(sql_str,language='cn')
# m_price.tail()
# #豆粕价格整体
# path=path_oil_meal+'\\豆粕现货价格.jpg'
# plot(m_price['2017':],'时间','价格','豆粕现货价格',path,marker='')
# #豆粕价格季节性以东莞为例
# m_price_dg=m_price['东莞豆粕价格']
# path=path_oil_meal+'\\东莞豆粕价格.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),m_price_dg,'2018','元/吨','东莞豆粕价格',path)
#
# #菜粕现货价格
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'rapeseed_meal_price%%'"""
# rm_price= fetch_edb(sql_str,language='cn')
# rm_price.tail()
# path=path_oil_meal+'\\菜粕现货价格.jpg'
# plot(rm_price.dropna()['2017':],'时间','价格','菜粕现货价格',path,marker='')
#
#
# #豆菜粕价差
# m_rm_nt=m_price['张家港豆粕价格']-rm_price['南通菜粕价格']
# m_rm_nt=m_rm_nt.resample('D').mean().interpolate()#差值处理
# path=path_oil_meal+'\\南通豆菜粕价差.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),m_rm_nt,'2018','价差','南通豆菜粕价差',path)
# #豆油现货价格
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('soybean_oil_price_average','soybean_oil_price_Huangpu','soybean_oil_price_Zhangjiagang',
# 'soybean_oil_price_Rizhao','soybean_oil_price_Tianjin','soybean_oil_price_Ningbo')"""
# y_price_1= fetch_edb(sql_str,language='cn')
# y_price_1.tail()
# #一级豆油价格走势
# path=path_oil_meal+'\\一级豆油价格.jpg'
# plot(y_price_1['2016':],'时间','价格','一级豆油现货价格',path,marker='')
#
# #豆油出厂价格
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('soybean_oil_price_fab_Fangchenggang','soybean_oil_price_fab_Zhanjiang','soybean_oil_price_fab_Zhangjiagang',
# 'soybean_oil_price_fab_Rizhao','soybean_oil_price_fab_Jingjin','soybean_oil_price_fab_Dalian')"""
# y_fab_price= fetch_edb(sql_str,language='cn')
# path=path_oil_meal+'\\豆油出厂价格.jpg'
# plot(y_fab_price['2016':],'时间','价格','豆油出厂价格',path,marker='')
#
# y1_price_zjg=y_price_1['张家港豆油价格']
# path=path_oil_meal+'\\张家港一级豆油价格季节性.jpg'
# period_plot(list(['2010','2011','2012','2013','2014','2015','2016','2017']),y1_price_zjg,'2018','元/吨','张家港一级豆油价格季节性',path)
#
#
# #菜油现货价格
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'grade4_rapeseed_oil_price%%'"""
# oi_price= fetch_edb(sql_str,language='cn')
# oi_price.tail()
# oi4_price_nt=oi_price['南通四级菜油出厂价格'].resample('D').mean().interpolate()#差值处理
# path=path_oil_meal+'\\南通四级菜籽油价格走势.jpg'
# period_plot(list(['2010','2011','2012','2013','2014','2015','2016','2017']),oi4_price_nt,'2018','元/吨','南通四级菜籽油价格走势',path)
#
# path=path_oil_meal+'\\四级菜籽油出厂价格.jpg'
# plot(oi_price['2016':],'时间','价格','四级菜籽油出厂价格',path,marker='')
#
#
#
# #棕榈现货价格
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where cn_name = '马来西亚本地天然棕榈油价格'"""
# p_price_Malaysia= fetch_edb(sql_str,language='cn')
# path=path_oil_meal+'\\马来西亚本地天然棕榈油价格.jpg'
# plot(p_price_Malaysia['2015':].resample('D').mean().interpolate(),'时间','林吉特/吨','马来西亚本地天然棕榈油价格',path,marker='')
#
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('palm_oil_price_average','palm_oil_price_average_24','palm_oil_price_Fujian','palm_oil_price_Guangdong'
# ,'palm_oil_price_Zhangjiagang','palm_oil_price_Tianjin','palm_oil_price_Rizhao')"""
# p_price= fetch_edb(sql_str,language='cn')
# p_price.tail()
# path=path_oil_meal+'\\棕榈油价格.jpg'
# plot(p_price['2015':].resample('D').mean().interpolate(),'时间','价格','棕榈油价格',path,marker='')
#
# p_price_hp=p_price['广东24度棕榈油价格'].resample('D').mean().interpolate()#差值处理
# path=path_oil_meal+'\\广东24度棕榈油价格.jpg'
# period_plot(list(['2010','2011','2012','2013','2014','2015','2016','2017']),p_price_hp,'2018','元/吨','广东24度棕榈油价格',path)
#
#
#
# #棕榈油进口贸易
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('palm_oil_fob_Malaysia','palm_oil_cnf_Malaysia','palm_oil_import_cost')"""
# p_price_fob= fetch_edb(sql_str,language='cn')
# p_price_fob.tail()
#
# path=path_oil_meal+'\\马来西亚棕榈油FOB及进口成本.jpg'
# twin_plot(p_price_fob[['棕榈油马来西亚到岸价','棕榈油马来西亚离岸价']],p_price_fob[['棕榈油马来西亚进口成本价']]
#           ,path,y_label1='美元/吨',y_label2='元/吨',title='马来西亚棕榈油FOB及进口成本')
#
# p_price_fob.tail()
#
# #棕榈油进口贸易
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where cn_name in ('进口棕榈油装船量当月值','进口马来西亚棕榈油装船量当月值','进口印度尼西亚棕榈油装船量当月值','进口棕榈油到港量当月值')"""
# p_shipment= fetch_edb(sql_str,language='cn')
# path=path_oil_meal+'\\进口棕榈油装船量当月值.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),p_shipment['进口棕榈油装船量当月值'].resample('M').sum().interpolate(),'2018','吨','进口棕榈油装船量当月值',path)
# path=path_oil_meal+'\\进口棕榈油到港量当月值.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),p_shipment['进口棕榈油到港量当月值'].resample('M').sum().interpolate(),'2018','吨','进口棕榈油装船量当月值',path)
#
# #油脂油料进口海关统计
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where cn_name in ('海关口径棕榈油进口量当月值','海关口径棕榈油进口平均价','海关口径棕榈油进口量累计值','海关口径棕榈油进口量累计同比')"""
# p_custom_statistics= fetch_edb(sql_str,language='cn')
# path=path_oil_meal+'\\海关口径棕榈油进口量累计值.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),p_custom_statistics['海关口径棕榈油进口量累计值'].interpolate(),'2018','万吨','海关口径棕榈油进口量累计值',path)
# path=path_oil_meal+'\\海关口径棕榈油进口量累计同比.jpg'
# plot(pd.DataFrame(p_custom_statistics['海关口径棕榈油进口量累计同比']['2013':],columns=['海关口径棕榈油进口量累计同比']),'时间','%','海关口径棕榈油进口量累计同比',path,marker='')
#
# path=path_oil_meal+'\\海关口径棕榈油进口量与进口成本.jpg'
# twin_plot(p_custom_statistics[['海关口径棕榈油进口量当月值']],p_custom_statistics[['海关口径棕榈油进口平均价']]
#           ,path,y_label1='万吨',y_label2='美元/吨',title='海关口径棕榈油进口量与进口成本')
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where cn_name in ('当月大豆进口数量','全国进口大豆船期预报到港量')"""
# soybean_custom_statistics= fetch_edb(sql_str,language='cn')
# path=path_oil_meal+'\\当月大豆进口数量.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),soybean_custom_statistics['当月大豆进口数量'],'2018','万吨','当月大豆进口数量',path)
# path=path_oil_meal+'\\全国进口大豆船期预报到港量.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),soybean_custom_statistics['全国进口大豆船期预报到港量'],'2018','万吨','全国进口大豆船期预报到港量',path)
#
#
# #进口大豆压榨利润
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'imported_soybeans_crush_margin%%'"""
# soybean_crush_margin=fetch_edb(sql_str,language='cn')
# soybean_crush_margin.tail()
#
# path=path_oil_meal+'\\大豆压榨利润.jpg'
# plot(soybean_crush_margin['2016':],'时间','价格','大豆压榨利润',path,marker='')
#
#
# #豆棕现货价差
# yp=pd.DataFrame((y_price_1['日照豆油价格']-p_price['日照24度棕榈油价格']).dropna(),columns=['日照豆棕价差'])
# path=path_oil_meal+'\\日照豆棕价差.jpg'
# plot(yp.dropna()['2017':],'时间','价格','日照豆棕价差',path,marker='')
#
# #巴西大豆fob运费cnf
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'Brazilian_soybean%%'"""
# soybean_br_cost= fetch_edb(sql_str,language='cn')
# soybean_br_cost.tail()
# #巴西升贴水
# path=path_oil_meal+'\\巴西大豆近月升贴水.jpg'
# plot(soybean_br_cost[['巴西大豆近月升贴水']].resample('D').mean().interpolate(),'时间','美分/蒲式耳','巴西升贴水',path,marker='')
#
# path=path_oil_meal+'\\巴西豆升贴水季节性.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),
#             soybean_br_cost['巴西大豆近月升贴水'].resample('D').mean().interpolate(),'2018','美分/蒲式耳','巴西豆升贴水季节性',path)
#
#
# # import_cost=soybean_br_cost.soybean_br_cnf*6.3*1.03*1.13+120
# # import_cost.tail()
# # #进口大豆到港成本计算公式=（CBOT大豆价格+升贴水）*单位换算0.367437*汇率6.83*关税1.03*增值税1.13+港杂费120
#
#
# soybean_br_cost_fob=soybean_br_cost[['巴西大豆近月CNF价格','巴西大豆近月FOB']]
# Brazilian_soybean_onshore_duty_price=soybean_br_cost[['巴西大豆近月运费']]
# path=path_oil_meal+'\\巴西大豆CNF和运费.jpg'
# twin_plot(soybean_br_cost_fob,Brazilian_soybean_onshore_duty_price['2015':],path,y_label1='美元/吨',y_label2='元/吨',title='巴西大豆CNF和运费')
#
# path=path_oil_meal+'\\巴西大豆完税到岸成本.jpg'
# plot(soybean_br_cost[['巴西大豆到岸完税价']].resample('D').mean().interpolate(),'时间','元/吨','巴西大豆完税到岸成本',path,marker='')
# #美国大豆运费cnf
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('cost_of_importing_American_soybean_into_China',
# 'freight_of_importing_American_soybean_into_China','American_soybean_premiums_discount_in_recent_month')"""
# soybean_us_cost= fetch_edb(sql_str)
# soybean_us_cost.tail()
# path=path_oil_meal+'\\美豆升贴水.jpg'
# period_plot(list(['2011','2012','2013','2014','2015','2016','2017']),
#             soybean_us_cost.American_soybean_premiums_discount_in_recent_month.resample('D').mean().interpolate(),'2018','美分/蒲式耳','美豆升贴水',path)
#
# fig = plt.figure()
# ax1 = fig.add_subplot(111)
# ax1.plot(soybean_us_cost.index,soybean_us_cost.American_soybean_premiums_discount_in_recent_month,label='American_soybean_premiums_discount_in_recent_month')
# ax1.plot(soybean_us_cost.index,soybean_us_cost.freight_of_importing_American_soybean_into_China,label='freight_of_importing_American_soybean_into_China')
# ax1.set_ylabel('美豆升贴水（美分/蒲）和运费（美元/吨）')
# ax1.legend(loc=1)
# ax2=ax1.twinx()
# ax2.plot(soybean_us_cost.index,soybean_us_cost.cost_of_importing_American_soybean_into_China,'r',label='cost_of_importing_American_soybean_into_China')
# ax2.legend(loc=2)
# ax2.set_ylabel('美豆到岸完税成本/RMB/吨')
# plt.title('美国大豆升贴水、运费和到岸完税成本',fontsize=16)
# fig = plt.gcf()
# fig.set_size_inches(8, 5)
# plt.savefig(r'D:\future\oil_meal\美国大豆进口成本.jpg',dpi=240)
# plt.close()
#
#
# #soybean_us_cost['2017-01-01':]['soybean_us_cost'].plot()
# #CBOT大豆价格
#
#
# #油粕比价
# ym_rz=y_price_1['日照豆油价格']/m_price['日照豆粕价格']
# ym_rz.tail()
#
# plt.plot(ym_rz['2015':].dropna())
# plt.title('日照一级豆油/豆粕')
# fig = plt.gcf()
# fig.set_size_inches(8, 5)
# plt.savefig(r'D:\future\oil_meal\日照一级豆油比豆粕.jpg',dpi=240)
# plt.close()
#
#
# oim_nt_zjg=oi_price['南通四级菜油出厂价格']/m_price['张家港豆粕价格']
# oim_nt_zjg.tail()
# plt.plot(oim_nt_zjg.dropna()['2015':])
# plt.title('南通四级菜油/日照豆粕')
# fig = plt.gcf()
# fig.set_size_inches(8, 5)
# plt.savefig(r'D:\future\oil_meal\南通四级菜油比日照豆粕.jpg',dpi=240)
# plt.close(1)
#
# #oim_nt_zjg['2017'].dropna().plot()
# oim_nt_zjg.dropna().mean()
# bigger_aver=oim_nt_zjg.dropna()>oim_nt_zjg.dropna().mean()
#
# """
# #M1805
# m1805= pd.DataFrame(list(db.futures.find({'beacon':'m1805'}))).dropna()
# m1805.index=pd.to_datetime(m1805.date)
# del m1805['date']
# m1805.CLOSE
# #base
# jicha_rz=m_price['2018'].m_price_rz-m1805.CLOSE['2018']
# jicha_rz.dropna().plot()
# jicha_dg=m_price['2018'].m_price_dg-m1805.CLOSE['2018']
# jicha_dg.dropna().plot()
# plt.close(1)
#
# y1805= pd.DataFrame(list(db.futures.find({'beacon':'y1805'}))).dropna()
# y1805.index=pd.to_datetime(y1805.date)
# del y1805['date']
#
# #组合的移动标准差
# oi1805= pd.DataFrame(list(db.futures.find({'beacon':'oi805'}))).dropna()
# oi1805.index=pd.to_datetime(oi1805.date)
# del oi1805['date']
# oim=oi1805.CLOSE-2*m1805.CLOSE
# oimstd=pd.rolling_std(oim,20)
# mstd=pd.rolling_std(m1805.CLOSE,20)
# oistd=pd.rolling_std(oi1805.CLOSE,20)
#
# plt.plot(oimstd,label='oimstd')
# plt.plot(mstd,label='mstd')
# plt.plot(oistd,label='oistd')
# plt.legend()
#
# oiy=oi1805.CLOSE-y1805.CLOSE
# oiystd=pd.rolling_std(oiy,20)
# ystd=pd.rolling_std(y1805.CLOSE,20)
# plt.plot(ystd,label='ystd')
# plt.plot(oiystd,label='oiystd')
# plt.close()
# """
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'soybean_meal_inventory%%'"""
# meal_inventory=fetch_edb(sql_str)
# path=path_oil_meal+'\\全国豆粕库存.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),meal_inventory['soybean_meal_inventory_of_China'],'2018','万吨','全国豆粕库存',path)
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in( 'soybean_consuming_of_China_port','soybean_inventory_of_China_port')"""
# soybean_inventory=fetch_edb(sql_str)
# path=path_oil_meal+'\\港口大豆消耗.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),(soybean_inventory['soybean_consuming_of_China_port']).rolling(window=5,center=False).mean(),'2018','吨','港口大豆消耗',path)
#
# path=path_oil_meal+'\\港口大豆库存.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),(soybean_inventory['soybean_inventory_of_China_port']).rolling(window=5,center=False).mean(),'2018','吨','港口大豆库存',path)
# #豆油商业库存
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'soybean_oil_business_inventory_of_China'"""
# business_inventory_y=fetch_edb(sql_str)
# business_inventory_y=business_inventory_y.resample('D').mean().interpolate()#差值处理
# path=path_oil_meal+'\\豆油商业库存.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),business_inventory_y,'2018','万吨','豆油商业库存',path)
# #棕榈油港口
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'palm_oil_inventory%%'"""
# port_inventory_p=fetch_edb(sql_str)
# path=path_oil_meal+'\\棕榈油港口库存.jpg'
# period_plot(list(['2011','2012','2013','2014','2015','2016','2017']),
#             port_inventory_p.palm_oil_inventory_of_China_port.resample('D').mean().interpolate(),'2018','万吨','棕榈油港口库存',path)
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where cn_name in ('马来西亚棕榈油期末库存','马来西亚棕榈仁油期末库存')"""
# palm_oil_inventory_of_Malaysia=fetch_edb(sql_str,language='cn')
#
#
# #棕榈油产量
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'palm_oil_production%%'"""
# p_production_my=fetch_edb(sql_str)
# path=path_oil_meal+'\\马来西亚棕榈油产量.jpg'
# period_plot(list(['2011','2012','2013','2014','2015','2016','2017']),
#             p_production_my.palm_oil_production_of_Malaysia/10000,'2018','万吨','马来西亚棕榈油产量',path)
# #外汇市场
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'Brazilian_real%%'"""
# brl=fetch_edb(sql_str)
# path=path_oil_meal+'\\美元兑巴西雷亚尔.jpg'
# plot(brl['2015':][['Brazilian_real_buy']],'时间','Brazilian_real_buy','美元兑巴西雷亚尔',path,marker='')
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'Argentine_peso'"""
# peso=fetch_edb(sql_str)
# path=path_oil_meal+'\\阿根廷比索.jpg'
# plot(peso['2016':],'时间','peso','阿根廷比索',path,marker='')
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('Malaysian_ringgit') """
# Malaysian_ringgit=fetch_edb(sql_str)
# path=path_oil_meal+'\\Malaysian_ringgit.jpg'
# plot(Malaysian_ringgit['2013':],'日期','Malaysian_ringgit','美元兑林吉特',path)
#
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name in ('Indonesia_rupiah_buy') """
# Indonesia_rupiah_buy=fetch_edb(sql_str)
# path=path_oil_meal+'\\Indonesia_rupiah_buy.jpg'
# plot(Indonesia_rupiah_buy['2013':],'日期','Indonesia_rupiah_buy','美元兑印尼卢比',path)
#
#
# #大豆进口预到港
# sql_str="""SELECT info.cn_name,info.en_name,edb.trade_date,edb.value FROM wind_commodity_edb edb left join wind_commodity_info info
# on edb.wind_code=info.wind_code where en_name like 'forcasting_of_soybean_importing'"""
# forecast_import_soy=fetch_edb(sql_str)
# #forecast_import_soy.nation.plot()
# path=path_oil_meal+'\\大豆进口预到港.jpg'
# period_plot(list(['2013','2014','2015','2016','2017']),
#            forecast_import_soy.forcasting_of_soybean_importing.resample('M').mean().interpolate(),'2018','万吨','全国大豆进口预到港',path)
#
# # #美豆出口中国检验量
# # us_soy_inspection_cn=fetch('us_soy_inspection_cn','oil_meal')*27.2/10000
# # path=path_oil_meal+'\\美豆出口中国检验量.jpg'
# # crop_plot(list(['2010','2011','2012','2013','2014','2015','2016','2017','2018']),
# #            us_soy_inspection_cn.us_soy_inspection_cn.resample('W').mean().interpolate(),'万吨','美豆出口中国检验量',path)
# #
# # crop_year_not_shipped_soy=fetch('crop_year_not_shipped_soy','oil_meal')
# # path=path_oil_meal+'\\当年市场年度美豆出口中国未装船量.jpg'
# # crop_plot(list(['2010','2011','2012','2013','2014','2015','2016','2017','2018']),
# #           crop_year_not_shipped_soy.crop_year_not_shipped_cn/10000,'万吨','当年市场年度美豆出口中国未装船量',path)

if __name__=="__main__":
    add_ppt_oil_meal()